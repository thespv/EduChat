import os
import base64
from pathlib import Path
from io import BytesIO
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from typing import Optional, List
from dotenv import load_dotenv

from api.database import (
    init_db, create_session, get_sessions, get_session, 
    add_message, update_session_title, delete_session,
    save_lecture_note, get_lecture_notes, delete_lecture_note
)

# Load .env from project root (works locally and on production)
env_path = Path(__file__).parent.parent / ".env"
load_dotenv(env_path)

GEMINI_KEY = os.getenv("GEMINI_API_KEY", "")
os.environ["GEMINI_API_KEY"] = GEMINI_KEY

app = FastAPI()

init_db()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Import both services
from api.services.llm_service import process_multimodal_query
from api.services.langchain_service import (
    process_rag_query, 
    generate_quiz_with_rag,
    generate_flashcards_with_rag,
    rag_system
)

GEMINI_KEY = os.getenv("GEMINI_API_KEY", "")
print(f"Index.py - Env path: {env_path}")
print(f"Index.py - GEMINI_KEY: {GEMINI_KEY[:20] if GEMINI_KEY else 'NOT FOUND'}")

os.environ["GEMINI_API_KEY"] = GEMINI_KEY

MAX_FILE_SIZE = 4.5 * 1024 * 1024

@app.post("/api/chat")
async def chat_endpoint(
    message: str = Form(...),
    user: str = Form(default="User"),
    session_id: Optional[int] = Form(default=None),
    save_history: bool = Form(default=True),
    files: Optional[List[UploadFile]] = File(default=None)
):
    try:
        processed_files = []
        
        if files:
            for file in files:
                content = await file.read()
                
                if len(content) > MAX_FILE_SIZE:
                    return JSONResponse({
                        "error": f"File '{file.filename}' exceeds 4.5MB limit. Please use smaller files."
                    }, status_code=413)
                
                file_type = file.content_type or "application/octet-stream"
                
                if file_type.startswith("image/"):
                    b64_data = base64.b64encode(content).decode("utf-8")
                    processed_files.append({
                        "type": "image",
                        "name": file.filename,
                        "data": b64_data,
                        "mime": file_type
                    })
                elif file_type.startswith("audio/"):
                    b64_data = base64.b64encode(content).decode("utf-8")
                    processed_files.append({
                        "type": "audio",
                        "name": file.filename,
                        "data": b64_data,
                        "mime": file_type
                    })
                elif file_type.startswith("video/"):
                    return JSONResponse({
                        "error": "Video processing requires external storage. For now, please provide a video URL or extract audio."
                    }, status_code=400)
                elif file.filename.endswith(".pdf"):
                    b64_data = base64.b64encode(content).decode("utf-8")
                    processed_files.append({
                        "type": "pdf",
                        "name": file.filename,
                        "data": b64_data
                    })
                elif file.filename.endswith((".txt", ".md")):
                    text_content = content.decode("utf-8", errors="ignore")
                    processed_files.append({
                        "type": "text",
                        "name": file.filename,
                        "data": text_content
                    })
                elif file.filename.endswith((".doc", ".docx")):
                    b64_data = base64.b64encode(content).decode("utf-8")
                    processed_files.append({
                        "type": "doc",
                        "name": file.filename,
                        "data": b64_data
                    })
                elif file.filename.endswith((".ppt", ".pptx")):
                    b64_data = base64.b64encode(content).decode("utf-8")
                    processed_files.append({
                        "type": "pptx",
                        "name": file.filename,
                        "data": b64_data
                    })
                else:
                    processed_files.append({
                        "type": "unknown",
                        "name": file.filename,
                        "data": content.decode("utf-8", errors="ignore")[:1000]
                    })
        
        # Use LangChain RAG if files attached, otherwise use regular LLM
        if processed_files:
            # RAG mode with document context
            try:
                result = await process_rag_query(message, user, processed_files, session_id or 0)
            except Exception as e:
                print(f"RAG error, falling back to regular LLM: {e}")
                result = await process_multimodal_query(message, user, processed_files)
        else:
            # Regular LLM mode
            result = await process_multimodal_query(message, user, processed_files)
        
        if save_history:
            if not session_id:
                session_id = create_session(user, message[:30] + "...")
            add_message(session_id, "user", message)
            add_message(session_id, "bot", result)
        
        return {"reply": result, "session_id": session_id}
        
    except Exception as e:
        return JSONResponse({
            "error": f"Error processing request: {str(e)}"
        }, status_code=500)

@app.get("/api/health")
async def health_check():
    return {"status": "healthy", "service": "EduChat API"}

@app.get("/api/chat/history")
async def get_chat_history(user: str = "User"):
    try:
        sessions = get_sessions(user)
        return {"sessions": sessions}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.get("/api/chat/session/{session_id}")
async def get_chat_session(session_id: int):
    try:
        session = get_session(session_id)
        if not session:
            return JSONResponse({"error": "Session not found"}, status_code=404)
        return session
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/api/chat/session")
async def create_chat_session(user: str = Form(...), title: str = Form(default="New Chat")):
    try:
        print(f"Creating session for user: {user}, title: {title}")
        session_id = create_session(user, title)
        print(f"Session created with ID: {session_id}")
        return {"session_id": session_id, "title": title}
    except Exception as e:
        print(f"Error creating session: {e}")
        return JSONResponse({"error": str(e)}, status_code=500)

@app.put("/api/chat/session/{session_id}")
async def update_chat_session(session_id: int, title: str = Form(...)):
    try:
        update_session_title(session_id, title)
        return {"success": True}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.delete("/api/chat/session/{session_id}")
async def delete_chat_session(session_id: int):
    try:
        delete_session(session_id)
        return {"success": True}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/api/extract-pdf")
async def extract_pdf(content: str = Form(...), file_type: str = Form(default="pdf")):
    try:
        import base64
        from io import BytesIO
        from fastapi.responses import StreamingResponse
        import json
        import asyncio
        
        if file_type == "pdf":
            from pypdf import PdfReader
            
            async def generate_pages():
                pdf_bytes = base64.b64decode(content)
                pdf_file = BytesIO(pdf_bytes)
                reader = PdfReader(pdf_file)
                total_pages = len(reader.pages)
                
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    
                    page_data = {
                        "page_num": i + 1,
                        "total_pages": total_pages,
                        "text": text[:8000],
                        "extracted": True
                    }
                    yield f"data: {json.dumps(page_data)}\n\n"
                    await asyncio.sleep(0.3)
                
                yield f"data: {json.dumps({'done': True})}\n\n"
            
            return StreamingResponse(generate_pages(), media_type="text/event-stream")
        
        elif file_type in ["doc", "docx"]:
            from docx import Document
            doc_bytes = base64.b64decode(content)
            doc_file = BytesIO(doc_bytes)
            doc = Document(doc_file)
            
            text = "\n".join([para.text for para in doc.paragraphs])
            return {"pages": [{"page_num": 1, "text": text}], "total_pages": 1, "success": True}
        
        elif file_type == "txt" or file_type == "md":
            text = base64.b64decode(content).decode('utf-8', errors='ignore')
            return {"pages": [{"page_num": 1, "text": text}], "total_pages": 1, "success": True}
        
        elif file_type in ["ppt", "pptx"]:
            try:
                from pptx import Presentation
                doc_bytes = base64.b64decode(content)
                doc_file = BytesIO(doc_bytes)
                prs = Presentation(doc_file)
                
                text = ""
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = f"\n--- Slide {slide_num} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text += shape.text + "\n"
                    text += slide_text
                
                return {"pages": [{"page_num": 1, "text": text[:50000]}], "total_pages": len(prs.slides), "success": True}
            except Exception as e:
                return {"error": f"Failed to extract PowerPoint: {str(e)}"}
        
        else:
            return {"error": "Unsupported format"}
            
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.get("/api/notes")
async def get_notes(user: str = "User"):
    try:
        notes = get_lecture_notes(user)
        return {"notes": notes}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/api/notes")
async def save_note(
    user: str = Form(...),
    name: str = Form(...),
    content: str = Form(...),
    file_type: str = Form(...)
):
    try:
        note_id = save_lecture_note(user, name, content, file_type)
        return {"success": True, "note_id": note_id}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.delete("/api/notes/{note_id}")
async def delete_note(note_id: int):
    try:
        delete_lecture_note(note_id)
        return {"success": True}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

# LangChain-powered endpoints
@app.post("/api/chat/generate-quiz")
async def generate_quiz(
    topic: str = Form(...),
    difficulty: str = Form(default="Medium"),
    num_questions: int = Form(default=5),
    user: str = Form(default="User"),
    use_rag: bool = Form(default=False)
):
    try:
        if use_rag:
            result = await generate_quiz_with_rag(topic, difficulty, num_questions, user, 0)
        else:
            from api.services.langchain_service import get_langchain_llm
            from langchain_core.prompts import ChatPromptTemplate
            
            llm = get_langchain_llm()
            prompt = ChatPromptTemplate.from_messages([
                ("system", f"You are a quiz generator. Create a {difficulty} quiz with {num_questions} multiple choice questions about the given topic.\nFormat each question as:\nQ1) [question]\nA) [option1]\nB) [option2]\nC) [option3]\nD) [option4]\nAnswer: [correct answer letter]"),
                ("human", "Topic: {topic}")
            ])
            chain = prompt | llm
            response = chain.invoke({"topic": topic})
            result = response.content if hasattr(response, 'content') else str(response)
        
        return {"quiz": result}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/api/chat/generate-flashcards")
async def generate_flashcards(
    topic: str = Form(...),
    num_cards: int = Form(default=5),
    user: str = Form(default="User"),
    use_rag: bool = Form(default=False)
):
    try:
        if use_rag:
            result = await generate_flashcards_with_rag(topic, num_cards, user, 0)
        else:
            from api.services.langchain_service import get_langchain_llm
            from langchain_core.prompts import ChatPromptTemplate
            
            llm = get_langchain_llm()
            prompt = ChatPromptTemplate.from_messages([
                ("system", f"You are an expert educator. Create {num_cards} flashcards about the given topic.\nFormat each card as:\nQ: [question]\nA: [answer]"),
                ("human", "Topic: {topic}")
            ])
            chain = prompt | llm
            response = chain.invoke({"topic": topic})
            result = response.content if hasattr(response, 'content') else str(response)
        
        return {"flashcards": result}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/api/rag/ingest")
async def ingest_document(
    file: UploadFile = File(...),
    doc_name: str = Form(default="document")
):
    try:
        content = await file.read()
        content_type = file.content_type or ""
        
        import base64
        b64_data = base64.b64encode(content).decode("utf-8")
        
        if content_type.startswith("image/"):
            return JSONResponse({"error": "Image RAG not supported yet"}, status_code=400)
        
        from api.services.langchain_service import rag_system
        
        if file.filename.endswith(".pdf"):
            pdf_bytes = base64.b64decode(b64_data)
            success = rag_system.process_pdf(pdf_bytes, doc_name)
        else:
            text = content.decode("utf-8", errors="ignore")
            success = rag_system.process_text(text, doc_name)
        
        if success:
            return {"success": True, "document": doc_name, "chunks": "created"}
        return JSONResponse({"error": "Failed to create embeddings"}, status_code=500)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.get("/api/rag/search")
async def search_documents(
    query: str,
    doc_name: Optional[str] = None,
    k: int = 3
):
    try:
        from api.services.langchain_service import rag_system
        results = rag_system.similarity_search(query, doc_name, k)
        return {"results": results}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)